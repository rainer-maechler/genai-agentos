#!/usr/bin/env python3
"""
Document Parser Agent for GenAI AgentOS Showcase

This agent handles various document formats and extracts raw content.
Capabilities:
- PDF text extraction
- Word document processing  
- Excel spreadsheet parsing
- PowerPoint slide extraction
- Image text recognition (OCR)
- Metadata extraction
"""

import asyncio
import json
import tempfile
import mimetypes
from pathlib import Path
from typing import Dict, Any, Optional, List
import aiofiles
from loguru import logger

# Document processing imports
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
from PIL import Image
import pytesseract

# GenAI Protocol
from genai_protocol import GenAISession


class DocumentParserAgent:
    """Agent that parses various document formats and extracts content."""
    
    def __init__(self):
        self.name = "document_parser"
        self.description = "Parses documents and extracts text content, metadata, and structure"
        self.version = "1.0.0"
        self.supported_formats = {
            'application/pdf': self._parse_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._parse_docx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._parse_xlsx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._parse_pptx,
            'text/plain': self._parse_text,
            'image/png': self._parse_image,
            'image/jpeg': self._parse_image,
            'image/jpg': self._parse_image,
        }
        
    async def process_message(self, session: GenAISession, message: str) -> Dict[str, Any]:
        """Process incoming message and parse document."""
        try:
            # Parse the input message
            request = json.loads(message)
            file_id = request.get('file_id')
            file_path = request.get('file_path')
            options = request.get('options', {})
            
            if not file_id and not file_path:
                return {
                    "error": "No file_id or file_path provided",
                    "status": "error"
                }
            
            # Get file content
            if file_id:
                file_content = await session.get_file_content(file_id)
                file_data = file_content.data
                mime_type = file_content.mime_type
                filename = file_content.filename
            else:
                # Load from file path
                async with aiofiles.open(file_path, 'rb') as f:
                    file_data = await f.read()
                mime_type, _ = mimetypes.guess_type(file_path)
                filename = Path(file_path).name
            
            # Parse document based on MIME type
            if mime_type not in self.supported_formats:
                return {
                    "error": f"Unsupported file format: {mime_type}",
                    "supported_formats": list(self.supported_formats.keys()),
                    "status": "error"
                }
            
            parser_func = self.supported_formats[mime_type]
            result = await parser_func(file_data, filename, options)
            
            # Add metadata
            result.update({
                "agent": self.name,
                "filename": filename,
                "mime_type": mime_type,
                "file_size": len(file_data),
                "status": "success"
            })
            
            logger.info(f"Successfully parsed {filename} ({mime_type})")
            return result
            
        except json.JSONDecodeError:
            return {"error": "Invalid JSON input", "status": "error"}
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            return {"error": str(e), "status": "error"}
    
    async def _parse_pdf(self, file_data: bytes, filename: str, options: Dict) -> Dict[str, Any]:
        """Parse PDF document."""
        try:
            with tempfile.NamedTemporaryFile(suffix='.pdf') as tmp_file:
                tmp_file.write(file_data)
                tmp_file.flush()
                
                with open(tmp_file.name, 'rb') as pdf_file:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    
                    # Extract metadata
                    metadata = pdf_reader.metadata or {}
                    
                    # Extract text from all pages
                    pages = []
                    full_text = ""
                    
                    for page_num, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        pages.append({
                            "page_number": page_num + 1,
                            "text": page_text,
                            "char_count": len(page_text)
                        })
                        full_text += page_text + "\n"
                    
                    return {
                        "document_type": "pdf",
                        "page_count": len(pdf_reader.pages),
                        "full_text": full_text.strip(),
                        "pages": pages,
                        "metadata": {
                            "title": metadata.get('/Title', ''),
                            "author": metadata.get('/Author', ''),
                            "subject": metadata.get('/Subject', ''),
                            "creator": metadata.get('/Creator', ''),
                            "producer": metadata.get('/Producer', ''),
                            "creation_date": str(metadata.get('/CreationDate', '')),
                            "modification_date": str(metadata.get('/ModDate', ''))
                        },
                        "word_count": len(full_text.split()),
                        "char_count": len(full_text)
                    }
                    
        except Exception as e:
            return {"error": f"PDF parsing failed: {str(e)}", "document_type": "pdf"}
    
    async def _parse_docx(self, file_data: bytes, filename: str, options: Dict) -> Dict[str, Any]:
        """Parse Word document."""
        try:
            with tempfile.NamedTemporaryFile(suffix='.docx') as tmp_file:
                tmp_file.write(file_data)
                tmp_file.flush()
                
                doc = Document(tmp_file.name)
                
                # Extract text from paragraphs
                paragraphs = []
                full_text = ""
                
                for para in doc.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        paragraphs.append({
                            "text": para_text,
                            "style": para.style.name if para.style else "Normal"
                        })
                        full_text += para_text + "\n"
                
                # Extract tables
                tables = []
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    tables.append(table_data)
                
                # Extract core properties
                props = doc.core_properties
                
                return {
                    "document_type": "docx",
                    "full_text": full_text.strip(),
                    "paragraphs": paragraphs,
                    "tables": tables,
                    "paragraph_count": len([p for p in paragraphs if p["text"]]),
                    "table_count": len(tables),
                    "metadata": {
                        "title": props.title or "",
                        "author": props.author or "",
                        "subject": props.subject or "",
                        "keywords": props.keywords or "",
                        "created": str(props.created) if props.created else "",
                        "modified": str(props.modified) if props.modified else "",
                        "last_modified_by": props.last_modified_by or ""
                    },
                    "word_count": len(full_text.split()),
                    "char_count": len(full_text)
                }
                
        except Exception as e:
            return {"error": f"DOCX parsing failed: {str(e)}", "document_type": "docx"}
    
    async def _parse_xlsx(self, file_data: bytes, filename: str, options: Dict) -> Dict[str, Any]:
        """Parse Excel spreadsheet."""
        try:
            with tempfile.NamedTemporaryFile(suffix='.xlsx') as tmp_file:
                tmp_file.write(file_data)
                tmp_file.flush()
                
                workbook = openpyxl.load_workbook(tmp_file.name, data_only=True)
                
                sheets = []
                full_text = ""
                
                for sheet_name in workbook.sheetnames:
                    worksheet = workbook[sheet_name]
                    
                    # Extract data from cells
                    sheet_data = []
                    sheet_text = ""
                    
                    for row in worksheet.iter_rows(values_only=True):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        if any(cell.strip() for cell in row_data):  # Skip empty rows
                            sheet_data.append(row_data)
                            sheet_text += " ".join(row_data) + "\n"
                    
                    sheets.append({
                        "name": sheet_name,
                        "data": sheet_data,
                        "row_count": len(sheet_data),
                        "col_count": worksheet.max_column,
                        "text": sheet_text.strip()
                    })
                    
                    full_text += f"Sheet: {sheet_name}\n{sheet_text}\n"
                
                return {
                    "document_type": "xlsx",
                    "full_text": full_text.strip(),
                    "sheets": sheets,
                    "sheet_count": len(sheets),
                    "metadata": {
                        "title": workbook.properties.title or "",
                        "author": workbook.properties.creator or "",
                        "subject": workbook.properties.subject or "",
                        "created": str(workbook.properties.created) if workbook.properties.created else "",
                        "modified": str(workbook.properties.modified) if workbook.properties.modified else ""
                    },
                    "word_count": len(full_text.split()),
                    "char_count": len(full_text)
                }
                
        except Exception as e:
            return {"error": f"XLSX parsing failed: {str(e)}", "document_type": "xlsx"}
    
    async def _parse_pptx(self, file_data: bytes, filename: str, options: Dict) -> Dict[str, Any]:
        """Parse PowerPoint presentation."""
        try:
            with tempfile.NamedTemporaryFile(suffix='.pptx') as tmp_file:
                tmp_file.write(file_data)
                tmp_file.flush()
                
                presentation = Presentation(tmp_file.name)
                
                slides = []
                full_text = ""
                
                for slide_num, slide in enumerate(presentation.slides):
                    slide_text = ""
                    shapes = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            shape_text = shape.text.strip()
                            shapes.append({
                                "text": shape_text,
                                "shape_type": str(shape.shape_type)
                            })
                            slide_text += shape_text + "\n"
                    
                    slides.append({
                        "slide_number": slide_num + 1,
                        "text": slide_text.strip(),
                        "shapes": shapes,
                        "shape_count": len(shapes)
                    })
                    
                    full_text += f"Slide {slide_num + 1}:\n{slide_text}\n"
                
                return {
                    "document_type": "pptx",
                    "full_text": full_text.strip(),
                    "slides": slides,
                    "slide_count": len(slides),
                    "metadata": {
                        "title": presentation.core_properties.title or "",
                        "author": presentation.core_properties.author or "",
                        "subject": presentation.core_properties.subject or "",
                        "created": str(presentation.core_properties.created) if presentation.core_properties.created else "",
                        "modified": str(presentation.core_properties.modified) if presentation.core_properties.modified else ""
                    },
                    "word_count": len(full_text.split()),
                    "char_count": len(full_text)
                }
                
        except Exception as e:
            return {"error": f"PPTX parsing failed: {str(e)}", "document_type": "pptx"}
    
    async def _parse_text(self, file_data: bytes, filename: str, options: Dict) -> Dict[str, Any]:
        """Parse plain text file."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text = None
            
            for encoding in encodings:
                try:
                    text = file_data.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                return {"error": "Unable to decode text file", "document_type": "text"}
            
            lines = text.split('\n')
            
            return {
                "document_type": "text",
                "full_text": text,
                "lines": [{"line_number": i+1, "text": line} for i, line in enumerate(lines)],
                "line_count": len(lines),
                "metadata": {
                    "encoding": "utf-8",
                    "has_bom": text.startswith('\ufeff')
                },
                "word_count": len(text.split()),
                "char_count": len(text)
            }
            
        except Exception as e:
            return {"error": f"Text parsing failed: {str(e)}", "document_type": "text"}
    
    async def _parse_image(self, file_data: bytes, filename: str, options: Dict) -> Dict[str, Any]:
        """Parse image and extract text using OCR."""
        try:
            with tempfile.NamedTemporaryFile() as tmp_file:
                tmp_file.write(file_data)
                tmp_file.flush()
                
                # Open image
                image = Image.open(tmp_file.name)
                
                # Extract text using OCR
                ocr_text = pytesseract.image_to_string(image)
                
                return {
                    "document_type": "image",
                    "full_text": ocr_text.strip(),
                    "ocr_confidence": "high" if len(ocr_text.strip()) > 10 else "low",
                    "image_info": {
                        "format": image.format,
                        "mode": image.mode,
                        "size": image.size,
                        "width": image.width,
                        "height": image.height
                    },
                    "metadata": {
                        "has_text": len(ocr_text.strip()) > 0,
                        "ocr_method": "tesseract"
                    },
                    "word_count": len(ocr_text.split()),
                    "char_count": len(ocr_text)
                }
                
        except Exception as e:
            return {"error": f"Image parsing failed: {str(e)}", "document_type": "image"}


async def main():
    """Main function to run the Document Parser Agent."""
    agent = DocumentParserAgent()
    
    # Connect to GenAI AgentOS
    session = GenAISession(
        ws_url="ws://localhost:8080/ws",
        agent_name=agent.name,
        agent_description=agent.description,
        agent_version=agent.version
    )
    
    logger.info(f"Starting {agent.name} agent...")
    
    try:
        await session.connect()
        logger.info("Connected to GenAI AgentOS")
        
        # Listen for messages
        async for message in session.listen():
            logger.info(f"Received message: {message.id}")
            
            try:
                result = await agent.process_message(session, message.content)
                await session.send_response(message.id, json.dumps(result))
                logger.info(f"Sent response for message: {message.id}")
                
            except Exception as e:
                error_response = {
                    "error": f"Processing failed: {str(e)}",
                    "status": "error",
                    "agent": agent.name
                }
                await session.send_response(message.id, json.dumps(error_response))
                logger.error(f"Error processing message {message.id}: {str(e)}")
                
    except KeyboardInterrupt:
        logger.info("Agent stopped by user")
    except Exception as e:
        logger.error(f"Agent error: {str(e)}")
    finally:
        await session.disconnect()
        logger.info("Agent disconnected")


if __name__ == "__main__":
    asyncio.run(main())